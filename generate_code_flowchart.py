"""
generate_flowchart.py
Generates a PowerPoint (.pptx) flowchart of the landgen code structure.

Run this script from the upper-level landgen directory (.../tools/landgen/) 
    and it will write landgen_flowchart.pptx in the same directory.
    You can specify a custom output path as an optional argument.

Usage:
    python generate_flowchart.py            # writes landgen_flowchart.pptx
    python generate_flowchart.py <out.pptx> # writes to a custom path

Requires:
    conda install python-pptx
"""

import sys
import argparse
import ast
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import qn, nsdecls
from lxml import etree

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
C_ENTRY    = RGBColor(0x1F, 0x49, 0x7D)   # dark blue  – entry point
C_CORE     = RGBColor(0x2E, 0x75, 0xB6)   # mid blue   – core orchestration
C_MODULE   = RGBColor(0x70, 0xAD, 0x47)   # green      – top-level modules
C_SUB      = RGBColor(0xFF, 0xC0, 0x00)   # amber      – sub-steps
C_IO       = RGBColor(0xED, 0x7D, 0x31)   # orange     – I/O utilities
C_DATA     = RGBColor(0x7B, 0x0E, 0xA8)   # purple     – shared data structures
C_PARALLEL = RGBColor(0xC0, 0x50, 0x4D)   # red        – parallel workers
C_NOTIMPL  = RGBColor(0xA5, 0xA5, 0xA5)   # grey       – not yet implemented
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK    = RGBColor(0x00, 0x00, 0x00)
C_ARROW    = RGBColor(0x40, 0x40, 0x40)

SLIDE_W = Inches(20)
SLIDE_H = Inches(14)

# Font sizes – adjust these to make all text larger or smaller
FS_TITLE = Pt(22)    # slide title
FS_LABEL = Pt(11)    # section labels / IO header
FS_BODY  = Pt(10)    # main box text
FS_SMALL = Pt(9.5)   # smaller boxes (substeps, workers, IO)
FS_ARROW = Pt(9)     # arrow labels

SCRIPT_DIR = Path(__file__).resolve().parent
SRC_DIR = SCRIPT_DIR / "src" / "landgen"
CONFIG_PATH = SCRIPT_DIR / "config.json"

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _rgb_to_hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def add_box(slide, left, top, width, height,
            text, fill_color, text_color=C_WHITE,
            font_size=FS_BODY, bold=False, shape_type="rect"):
    """Add a rounded-rectangle shape and return it."""

    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )

    # solid fill via python-pptx API (avoids duplicate prstGeom XML corruption)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color

    # thin white border
    shape.line.color.rgb = C_WHITE
    shape.line.width = Pt(0.75)

    # text
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.color.rgb = text_color

    return shape


def add_arrow(slide, x1, y1, x2, y2, color=C_ARROW, label=None):
    """Draw a straight connector with an arrowhead at (x2,y2).

    Builds the <p:cxnSp> from a complete XML template so the element has
    no <p:style> block (which can trigger PowerPoint repair on blank
    presentations) and no python-pptx connector-API side-effects.
    """
    # Bounding box of the line segment
    left = int(min(x1, x2))
    top  = int(min(y1, y2))
    cx   = int(abs(x2 - x1)) or 1   # OOXML requires cx >= 1
    cy   = int(abs(y2 - y1)) or 1

    # flip attributes encode which diagonal the line runs along
    flip_attrs = ''
    if x2 < x1:
        flip_attrs += ' flipH="1"'
    if y2 < y1:
        flip_attrs += ' flipV="1"'

    hex_color = _rgb_to_hex(color)
    sp_tree   = slide.shapes._spTree
    shape_id  = sp_tree._next_shape_id

    xml = (
        f'<p:cxnSp {nsdecls("a", "p", "r")}>'
        f'<p:nvCxnSpPr>'
        f'<p:cNvPr id="{shape_id}" name="Arrow {shape_id}"/>'
        f'<p:cNvCxnSpPr/>'
        f'<p:nvPr/>'
        f'</p:nvCxnSpPr>'
        f'<p:spPr>'
        f'<a:xfrm{flip_attrs}><a:off x="{left}" y="{top}"/>'
        f'<a:ext cx="{cx}" cy="{cy}"/></a:xfrm>'
        f'<a:prstGeom prst="line"><a:avLst/></a:prstGeom>'
        f'<a:ln w="19050">'
        f'<a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
        f'<a:tailEnd type="arrow" w="med" len="med"/>'
        f'</a:ln>'
        f'</p:spPr>'
        f'</p:cxnSp>'
    )

    cxnSp = parse_xml(xml)
    sp_tree.append(cxnSp)

    if label:
        mx = (x1 + x2) / 2
        my = (y1 + y2) / 2
        tb = slide.shapes.add_textbox(mx, my - Inches(0.12), Inches(0.9), Inches(0.2))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = label
        run.font.size = FS_ARROW
        run.font.color.rgb = C_ARROW

    return cxnSp


def _read_ast(py_path: Path):
    """Parse a Python file and return its AST, or None if unavailable."""
    try:
        return ast.parse(py_path.read_text(), filename=str(py_path))
    except Exception:
        return None


def _module_path(module_name: str) -> Path:
    return SRC_DIR / f"{module_name}.py"


def _module_exists(module_name: str) -> bool:
    return _module_path(module_name).exists()


def _top_level_functions(py_path: Path):
    """Return top-level function names from a Python file."""
    tree = _read_ast(py_path)
    if tree is None:
        return []
    return [n.name for n in tree.body if isinstance(n, ast.FunctionDef)]


def _discover_top_modules():
    """Read top-level module names from config.json modules list."""
    if CONFIG_PATH.exists():
        try:
            cfg = json.loads(CONFIG_PATH.read_text())
            names = [m.get("name") for m in cfg.get("modules", []) if m.get("name")]
            if names:
                return names
        except Exception:
            pass

    # fallback if config is unavailable
    fallback = []
    for p in sorted(SRC_DIR.glob("*.py")):
        stem = p.stem
        if stem.startswith("_"):
            continue
        if stem in {"__init__", "__main__", "tools", "shared_data", "landgen_io", "plot_landgen"}:
            continue
        fallback.append(stem)
    return fallback


def _discover_land_type_submodules():
    """Parse land_type.py for importlib.import_module('landgen.<module>') calls."""
    land_type_path = _module_path("land_type")
    tree = _read_ast(land_type_path)
    if tree is None:
        return []

    names = []
    for node in ast.walk(tree):
        if not isinstance(node, ast.Call):
            continue
        fn = node.func
        if not (isinstance(fn, ast.Attribute) and fn.attr == "import_module"):
            continue
        if not node.args:
            continue
        arg0 = node.args[0]
        if isinstance(arg0, ast.Constant) and isinstance(arg0.value, str):
            mod = arg0.value
            if mod.startswith("landgen."):
                short = mod.split(".")[-1]
                if short not in names:
                    names.append(short)
    return names


def _module_has_parallel_work(module_name: str) -> bool:
    """Heuristic: module uses mp.Pool or defines *_process worker functions."""
    py_path = _module_path(module_name)
    tree = _read_ast(py_path)
    if tree is None:
        return False

    for node in ast.walk(tree):
        if isinstance(node, ast.FunctionDef) and node.name.endswith("_process"):
            return True
        if isinstance(node, ast.Call):
            fn = node.func
            if isinstance(fn, ast.Attribute) and fn.attr == "Pool":
                return True
    return False


def _discover_landcover_deps():
    """Return relative-imported module names from landcover.py."""
    landcover_path = _module_path("landcover")
    tree = _read_ast(landcover_path)
    if tree is None:
        return []

    deps = []
    for node in ast.walk(tree):
        if isinstance(node, ast.ImportFrom) and node.level >= 1:
            for alias in node.names:
                name = alias.name
                if name not in deps and _module_exists(name):
                    deps.append(name)
    return deps


def _discover_io_items():
    """Build IO row items from parsed source functions."""
    items = []

    # Show remote-sensing helper when present.
    rs_path = _module_path("landcover_remote_sensing")
    if rs_path.exists():
        rs_funcs = [f for f in _top_level_functions(rs_path) if not f.startswith("_")]
        if rs_funcs:
            items.append(
                "landcover_remote_\nsensing.py\n" + " / ".join(f"{f}()" for f in rs_funcs[:2])
            )

    io_path = _module_path("landgen_io")
    io_funcs = [
        f for f in _top_level_functions(io_path)
        if f.startswith(("read_", "write_", "regrid_", "load_", "set_decomp_"))
    ]
    for f in io_funcs[:6]:
        items.append(f"{f}()")

    return items


def _brief_doc(node):
    """Return a short single-line summary from a function docstring."""
    doc = ast.get_docstring(node)
    if not doc:
        return "No docstring"
    first = doc.strip().splitlines()[0].strip()
    return first[:110]


def _function_summaries(module_name: str):
    """Return [(func_name, brief_doc), ...] for top-level functions in a module."""
    py_path = _module_path(module_name)
    tree = _read_ast(py_path)
    if tree is None:
        return []

    out = []
    for node in tree.body:
        if isinstance(node, ast.FunctionDef) and not node.name.startswith("_"):
            out.append((node.name, _brief_doc(node)))
    return out


def _discover_landgen_import_aliases(module_name: str, tree):
    """Map local import aliases to landgen module names for a module AST."""
    alias_to_module = {}

    for node in tree.body:
        if isinstance(node, ast.ImportFrom):
            # from . import landgen_io as io
            if node.level >= 1 and node.module is None:
                for alias in node.names:
                    target = alias.name
                    if _module_exists(target):
                        alias_to_module[alias.asname or alias.name] = target
            # from .landcover_remote_sensing import foo  (function/class import)
            # skip these for module-call matching.

        elif isinstance(node, ast.Import):
            # import landgen.landcover_remote_sensing as lc_rs
            for alias in node.names:
                name = alias.name
                if name.startswith("landgen."):
                    target = name.split(".")[-1]
                    if _module_exists(target):
                        alias_to_module[alias.asname or target] = target

    # Never treat self-module calls as cross-module.
    alias_to_module = {k: v for k, v in alias_to_module.items() if v != module_name}
    return alias_to_module


def _discover_process_calls(module_name: str):
    """Return all calls from *_process to functions in other landgen modules."""
    py_path = _module_path(module_name)
    tree = _read_ast(py_path)
    if tree is None:
        return []

    alias_to_module = _discover_landgen_import_aliases(module_name, tree)
    module_funcs_cache = {
        mod: set(_top_level_functions(_module_path(mod)))
        for mod in set(alias_to_module.values())
    }

    process_funcs = [
        n for n in tree.body
        if isinstance(n, ast.FunctionDef) and n.name.endswith("_process")
    ]
    if not process_funcs:
        return []

    calls = []
    for fn_node in process_funcs:
        for node in ast.walk(fn_node):
            if not isinstance(node, ast.Call):
                continue
            # Only include calls like <alias>.<func>() where alias is an imported
            # landgen module and func is defined in that target module.
            if not isinstance(node.func, ast.Attribute):
                continue
            if not isinstance(node.func.value, ast.Name):
                continue
            base = node.func.value.id
            func = node.func.attr
            if base not in alias_to_module:
                continue
            if func.startswith("_"):
                continue

            target_module = alias_to_module[base]
            target_funcs = module_funcs_cache.get(target_module, set())
            if func not in target_funcs:
                continue

            name = f"{base}.{func}"
            if name not in calls:
                calls.append(name)

    return calls


def build_functions_slide(prs, module_name: str, title_color=C_IO):
    """Add a slide listing top-level functions and brief descriptions."""
    slide_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(slide_layout)

    func_fs = Pt(18)

    # title
    tb = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(19.6), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f"{module_name}.py – Function Summary"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = title_color

    # subtitle
    st = slide.shapes.add_textbox(Inches(0.3), Inches(0.65), Inches(19.2), Inches(0.35))
    sp = st.text_frame.paragraphs[0]
    sr = sp.add_run()
    sr.text = f"Auto-parsed from src/landgen/{module_name}.py"
    sr.font.size = Pt(18)
    sr.font.color.rgb = C_BLACK

    funcs = _function_summaries(module_name)
    if not funcs:
        b = add_box(slide, Inches(0.5), Inches(1.3), Inches(19.0), Inches(0.8),
                    f"No top-level functions found in {module_name}.py", C_NOTIMPL,
                    text_color=C_WHITE, font_size=FS_BODY)
        return slide

    lines = [f"{name}()  –  {desc}" for name, desc in funcs]

    # Two-column list layout
    mid = (len(lines) + 1) // 2
    cols = [lines[:mid], lines[mid:]]
    lefts = [Inches(0.4), Inches(10.2)]
    top = Inches(1.1)
    width = Inches(9.4)
    height = Inches(12.4)

    for i, col_lines in enumerate(cols):
        box = slide.shapes.add_textbox(lefts[i], top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.clear()
        for j, line in enumerate(col_lines):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = line
            p.level = 0
            for r in p.runs:
                r.font.size = func_fs

    return slide


def build_overview_index_slide(prs):
    """Add an index slide summarizing discovered structure and where details live."""
    slide_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(slide_layout)

    # title
    tb = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(19.6), Inches(0.5))
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "landgen – Overview Index"
    run.font.size = FS_TITLE
    run.font.bold = True
    run.font.color.rgb = C_ENTRY

    top_modules = _discover_top_modules()
    submods = _discover_land_type_submodules()
    parallel = [m for m in submods if _module_exists(m) and _module_has_parallel_work(m)]

    io_funcs = [name for name, _ in _function_summaries("landgen_io")][:12]
    tools_funcs = [name for name, _ in _function_summaries("tools")][:10]
    plot_funcs = [name for name, _ in _function_summaries("plot_landgen")][:10]

    sections = [
        (
            "Execution Flow",
            "Entry: __main__.py -> landgen.main()\n"
            "Top-level modules: " + (", ".join(top_modules) if top_modules else "none"),
            C_CORE,
        ),
        (
            "land_type Submodules",
            "_process_single_year imports: " + (", ".join(submods) if submods else "none") + "\n"
            "Parallel workers: " + (", ".join(parallel) if parallel else "none"),
            C_SUB,
        ),
        (
            "landgen_io.py (Slide 3)",
            "Functions: " + (", ".join(io_funcs) if io_funcs else "none"),
            C_IO,
        ),
        (
            "tools.py (Slide 4) and plot_landgen.py (Slide 5)",
            "tools.py: " + (", ".join(tools_funcs) if tools_funcs else "none") + "\n"
            "plot_landgen.py: " + (", ".join(plot_funcs) if plot_funcs else "none"),
            C_ENTRY,
        ),
    ]

    y = Inches(0.9)
    h = Inches(2.85)
    for title, body, color in sections:
        add_box(
            slide,
            Inches(0.5),
            y,
            Inches(19.0),
            h,
            f"{title}\n{body}",
            color,
            text_color=C_WHITE,
            font_size=FS_SMALL,
        )
        y += h + Inches(0.2)

    return slide


# ---------------------------------------------------------------------------
# Layout constants  (all in Inches, converted via Inches())
# ---------------------------------------------------------------------------

def build_slide(prs):
    slide_layout = prs.slide_layouts[6]   # blank
    slide = prs.slides.add_slide(slide_layout)

    slide_fs = Pt(18)

    # ---- title ----
    tb = slide.shapes.add_textbox(Inches(0.2), Inches(0.1),
                                  Inches(19.6), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "landgen – Code Structure Flowchart"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = C_ENTRY

    # ---- legend ----
    legend_items = [
        (C_ENTRY,    "Entry point"),
        (C_CORE,     "Core orchestration"),
        (C_MODULE,   "Top-level module"),
        (C_SUB,      "Sub-step (sequential)"),
        (C_PARALLEL, "Parallel worker"),
        (C_IO,       "I/O utility"),
        (C_DATA,     "Shared data structure"),
        (C_NOTIMPL,  "Not yet implemented"),
    ]
    lx = Inches(0.2)
    ly = Inches(0.7)
    for i, (col, lbl) in enumerate(legend_items):
        add_box(slide, lx + Inches(i * 2.4), ly, Inches(2.2), Inches(0.55),
                lbl, col, font_size=slide_fs)

    # =========================================================
    # ROW 1 – Entry point
    # =========================================================
    r1y = Inches(1.5)
    bh = Inches(0.95)   # box height
    bw_sm  = Inches(3.1)
    bw_med = Inches(4.3)
    bw_lg  = Inches(3.6)

    # __main__.py
    main_x = Inches(7.9)
    b_main = add_box(slide, main_x, r1y, bw_med, bh,
                     "__main__.py\npython -m landgen <config.json>",
                     C_ENTRY, font_size=slide_fs, bold=True)

    # config.json (to the right)
    cfg_x = Inches(13.3)
    b_cfg = add_box(slide, cfg_x, r1y, bw_sm, bh,
                    "config.json\n(start/end year, paths, modules)",
                    C_IO, font_size=slide_fs)

    add_arrow(slide,
              cfg_x, r1y + bh / 2,
              main_x + bw_med, r1y + bh / 2)

    # =========================================================
    # ROW 2 – landgen.main()
    # =========================================================
    r2y = r1y + bh + Inches(0.65)
    core_x = Inches(7.3)
    bw_core = Inches(5.2)
    bh_core = Inches(1.6)
    b_core = add_box(slide, core_x, r2y, bw_core, bh_core,
                     "landgen.py  ·  main(config_path)\n"
                     "Load config · start mp.Manager · create GridData\n"
                     "Loop modules → importlib.import_module(name).run()",
                     C_CORE, font_size=slide_fs, bold=True)

    # arrow __main__ → main
    add_arrow(slide,
              main_x + bw_med / 2, r1y + bh,
              core_x + bw_core / 2, r2y)

    # =========================================================
    # ROW 3 – Shared data (right side)
    # =========================================================
    r3y = r2y
    sd_x = Inches(13.0)
    bw_sd = Inches(6.6)
    b_sd = add_box(slide, sd_x, r3y, bw_sd, Inches(2.15),
                   "shared_data.py  –  Multiprocessing shared structures\n"
                   "GridData / GridManager  (cell ids, lon/lat, landfrac)\n"
                   "TopoData / TopoManager  (elevation, slope, fmax …)\n"
                   "LtData  / LtManager     (pct_pft, harvest_frac, …)",
                   C_DATA, font_size=slide_fs)

    add_arrow(slide,
              core_x + bw_core, r2y + bh_core / 2,
              sd_x, r3y + Inches(1.075))

    # =========================================================
    # ROW 4 – Top-level modules  (5 boxes across)
    # =========================================================
    r4y = r3y + Inches(2.35) + Inches(0.75)
    top_modules = _discover_top_modules()
    modules = []
    for name in top_modules:
        col = C_MODULE if _module_exists(name) else C_NOTIMPL
        funcs = _top_level_functions(_module_path(name)) if _module_exists(name) else []
        text = f"{name}.py\nrun()"
        if name == "topography":
            text += "\n[must run first;\nsets landfrac]"
        elif name == "land_type":
            text += "\n[main land-type\norchestrator]"
        elif funcs:
            extras = [f for f in funcs if f not in {"run"} and not f.startswith("_")]
            if extras:
                text += "\n" + ", ".join(extras[:2])
        modules.append((text, col))

    if not modules:
        modules = [("No modules\nfound", C_NOTIMPL)]
    mod_bw   = Inches(3.6)
    mod_gap  = Inches(0.35)
    total_mw = len(modules) * mod_bw + (len(modules) - 1) * mod_gap
    mod_x0   = (SLIDE_W - total_mw) / 2
    mod_bh   = Inches(1.25)

    mod_boxes = []
    for i, (txt, col) in enumerate(modules):
        bx = mod_x0 + i * (mod_bw + mod_gap)
        b = add_box(slide, bx, r4y, mod_bw, mod_bh, txt, col,
                    font_size=slide_fs)
        mod_boxes.append((bx, b))
        # Branch arrows should originate from the bottom edge of landgen.main box.
        add_arrow(slide,
              core_x + bw_core / 2, r2y + bh_core,
              bx + mod_bw / 2, r4y)

    # =========================================================
    # ROW 5 – land_type internals: _process_single_year  loop
    # =========================================================
    r5y = r4y + mod_bh + Inches(0.8)
    land_type_idx = 0
    for i, name in enumerate(top_modules):
        if name == "land_type":
            land_type_idx = i
            break
    lt_x  = mod_x0 + land_type_idx * (mod_bw + mod_gap)
    bw_lt = Inches(3.4)
    bh_lt = Inches(0.9)
    b_loop = add_box(slide, lt_x, r5y, bw_lt, bh_lt,
                     "Loop years (start_year → end_year)\n"
                     "_process_single_year(year, …)",
                     C_CORE, font_size=slide_fs)
    add_arrow(slide,
              lt_x + bw_lt / 2, r4y + mod_bh,
              lt_x + bw_lt / 2, r5y)

    # =========================================================
    # ROW 6 – Sub-steps inside _process_single_year
    # =========================================================
    r6y = r5y + bh_lt + Inches(0.8)
    submod_names = _discover_land_type_submodules()
    substeps = []
    for name in submod_names:
        exists = _module_exists(name)
        parallel = _module_has_parallel_work(name) if exists else False
        col = C_SUB if exists else C_NOTIMPL
        txt = f"{name}.py\nrun()"
        if parallel:
            txt += " [parallel]"
        if not exists:
            txt += "\n[not yet impl.]"
        substeps.append((txt, col, name, parallel, exists))

    if not substeps:
        substeps = [("No sub-steps\nfound", C_NOTIMPL, "", False, False)]

    # Single-row substep layout: compact boxes so all modules fit on one row.
    ss_cols = max(1, len(substeps))
    ss_bw  = Inches(2.05)
    ss_bh  = Inches(1.0)
    ss_gap_x = Inches(0.12)
    ss_gap_y = Inches(0.0)
    row_w = ss_cols * ss_bw + (ss_cols - 1) * ss_gap_x
    ss_x0 = (SLIDE_W - row_w) / 2

    substep_centers = {}

    for i, (txt, col, name, parallel, exists) in enumerate(substeps):
        row = i // ss_cols
        col_i = i % ss_cols
        bx = ss_x0 + col_i * (ss_bw + ss_gap_x)
        by = r6y + row * (ss_bh + ss_gap_y)
        add_box(slide, bx, by, ss_bw, ss_bh, txt, col, font_size=slide_fs)
        substep_centers[name] = (bx + ss_bw / 2, by + ss_bh)
        add_arrow(slide,
                  lt_x + bw_lt / 2, r5y + bh_lt,
                  bx + ss_bw / 2, by)

    # =========================================================
    # ROW 7 – Parallel worker call detail
    # =========================================================
    n_sub_rows = (len(substeps) + ss_cols - 1) // ss_cols
    r7y = r6y + n_sub_rows * ss_bh + max(0, n_sub_rows - 1) * ss_gap_y + Inches(0.55)

    worker_modules = [name for _, _, name, parallel, exists in substeps if parallel and exists]
    if worker_modules:
        w_cols = min(3, len(worker_modules))
        w_bw = Inches(5.9)
        # Auto-size worker boxes so all discovered call lines fit.
        shown_worker_modules = worker_modules[:3]
        max_lines = 1
        for mod_name in shown_worker_modules:
            calls = _discover_process_calls(mod_name)
            # Account for visual wrapping of long call names.
            wrapped_call_lines = sum(max(1, (len(c) + 33) // 34) for c in calls) if calls else 1
            n_lines = 1 + wrapped_call_lines
            if n_lines > max_lines:
                max_lines = n_lines
        worker_line_h = 0.30
        worker_pad_h = 0.62
        w_bh = Inches(max(2.0, worker_pad_h + max_lines * worker_line_h))
        w_gap = Inches(0.35)
        total_ww = w_cols * w_bw + (w_cols - 1) * w_gap
        w_x0 = (SLIDE_W - total_ww) / 2

        for i, mod_name in enumerate(shown_worker_modules):
            wx = w_x0 + i * (w_bw + w_gap)
            proc_name = f"{mod_name}_process"
            calls = _discover_process_calls(mod_name)
            call_lines = "\n".join([f"· {c}()" for c in calls]) if calls else "· parallel chunk work"
            txt = f"{proc_name}()\n{call_lines}"

            add_box(slide, wx, r7y, w_bw, w_bh, txt, C_PARALLEL, font_size=slide_fs)

            if mod_name in substep_centers:
                sx, sy = substep_centers[mod_name]
                add_arrow(slide, sx, sy, wx + w_bw / 2, r7y)

    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        prog="generate_code_flowchart.py",
        description=(
            "Generate a PowerPoint (.pptx) flowchart of the landgen code structure."
        ),
        epilog=(
            "Examples:\n"
            "  python generate_code_flowchart.py\n"
            "  python generate_code_flowchart.py custom_flowchart.pptx"
        ),
        formatter_class=argparse.RawTextHelpFormatter,
    )
    parser.add_argument(
        "out_path",
        nargs="?",
        default="landgen_flowchart.pptx",
        help="Output .pptx file path (default: landgen_flowchart.pptx)",
    )
    args = parser.parse_args()

    out_path = Path(args.out_path)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    build_slide(prs)
    build_functions_slide(prs, "landgen_io", title_color=C_IO)
    build_functions_slide(prs, "tools", title_color=C_CORE)
    build_functions_slide(prs, "plot_landgen", title_color=C_ENTRY)

    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()

